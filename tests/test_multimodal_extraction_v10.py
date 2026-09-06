from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from app.multimodal_extraction_v10 import MultimodalExtractor, WorkspaceAssetResolver
from app.unified_multimodal_runtime_v10 import EvidenceKind, Modality, MultimodalRuntimeError


def test_workspace_resolver_rejects_escape(tmp_path: Path):
    outside = tmp_path.parent / "outside.txt"
    outside.write_text("nope", encoding="utf-8")
    resolver = WorkspaceAssetResolver(tmp_path)
    with pytest.raises(MultimodalRuntimeError):
        resolver.resolve(str(outside))


def test_text_file_extracts_native_evidence(tmp_path: Path):
    path = tmp_path / "notes.txt"
    path.write_text("DPN multimodal text", encoding="utf-8")
    result = MultimodalExtractor(tmp_path).extract("notes.txt", asset_id="notes")[0]
    assert result.asset.modality == Modality.TEXT
    assert result.asset.sha256
    assert result.evidence[0].kind == EvidenceKind.NATIVE_TEXT
    assert "multimodal" in result.evidence[0].content


def test_code_file_extracts_code_evidence(tmp_path: Path):
    path = tmp_path / "sample.py"
    path.write_text("print('dpn')", encoding="utf-8")
    result = MultimodalExtractor(tmp_path).extract("sample.py", asset_id="code")[0]
    assert result.asset.modality == Modality.CODE
    assert result.evidence[0].kind == EvidenceKind.CODE


def test_csv_becomes_spreadsheet_table_evidence(tmp_path: Path):
    path = tmp_path / "data.csv"
    path.write_text("name,value\nDPN,10\n", encoding="utf-8")
    result = MultimodalExtractor(tmp_path).extract("data.csv", asset_id="csv")[0]
    assert result.asset.modality == Modality.SPREADSHEET
    assert result.evidence[0].kind == EvidenceKind.TABLE
    assert "DPN | 10" in result.evidence[0].content


def test_docx_extracts_paragraphs_and_tables(tmp_path: Path):
    path = tmp_path / "report.docx"
    doc = Document()
    doc.add_paragraph("Executive summary")
    table = doc.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    doc.save(path)

    result = MultimodalExtractor(tmp_path).extract("report.docx", asset_id="doc")[0]
    kinds = {item.kind for item in result.evidence}
    assert result.asset.modality == Modality.DOCUMENT
    assert EvidenceKind.NATIVE_TEXT in kinds
    assert EvidenceKind.TABLE in kinds


def test_xlsx_extracts_structure_and_table(tmp_path: Path):
    path = tmp_path / "book.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Summary"
    ws.append(["Metric", "Value"])
    ws.append(["Sales", 42])
    wb.save(path)

    result = MultimodalExtractor(tmp_path).extract("book.xlsx", asset_id="sheet")[0]
    kinds = {item.kind for item in result.evidence}
    assert result.asset.modality == Modality.SPREADSHEET
    assert kinds == {EvidenceKind.STRUCTURE, EvidenceKind.TABLE}
    assert "Summary" in result.evidence[0].content


def test_pptx_extracts_slide_structure_and_text(tmp_path: Path):
    path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "DPN AI"
    slide.placeholders[1].text = "Unified multimodal intelligence"
    prs.save(path)

    result = MultimodalExtractor(tmp_path).extract("deck.pptx", asset_id="deck")[0]
    kinds = {item.kind for item in result.evidence}
    assert result.asset.modality == Modality.PRESENTATION
    assert EvidenceKind.STRUCTURE in kinds
    assert EvidenceKind.NATIVE_TEXT in kinds
    assert "DPN AI" in result.asset.native_text


def test_image_inventory_does_not_fabricate_visual_evidence(tmp_path: Path):
    path = tmp_path / "image.png"
    path.write_bytes(b"\x89PNG\r\n\x1a\nnot-a-real-image-but-provenance-only")
    result = MultimodalExtractor(tmp_path).extract("image.png", asset_id="image")[0]
    assert result.asset.modality == Modality.IMAGE
    assert [item.kind for item in result.evidence] == [EvidenceKind.METADATA]


def test_audio_inventory_does_not_claim_transcript(tmp_path: Path):
    path = tmp_path / "sample.wav"
    path.write_bytes(b"RIFFfake")
    result = MultimodalExtractor(tmp_path).extract("sample.wav", asset_id="audio")[0]
    assert result.asset.modality == Modality.AUDIO
    assert all(item.kind != EvidenceKind.TRANSCRIPT for item in result.evidence)


def test_attach_transcript_requires_audio_or_video_and_provenance(tmp_path: Path):
    path = tmp_path / "sample.wav"
    path.write_bytes(b"RIFFfake")
    extracted = MultimodalExtractor(tmp_path).extract("sample.wav", asset_id="audio")[0]
    evidence = MultimodalExtractor.attach_transcript(
        extracted.asset,
        "spoken words",
        provider="whisper",
        model="large-v3",
    )
    assert evidence.kind == EvidenceKind.TRANSCRIPT
    assert evidence.provider == "whisper"


def test_unknown_extension_fails_closed(tmp_path: Path):
    path = tmp_path / "unknown.bin"
    path.write_bytes(b"binary")
    with pytest.raises(MultimodalRuntimeError):
        MultimodalExtractor(tmp_path).extract("unknown.bin")
