from __future__ import annotations

import csv
import hashlib
import json
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from app.unified_multimodal_runtime_v10 import EvidenceKind, Modality, MultimodalAsset, MultimodalEvidence, MultimodalRuntimeError


TEXT_SUFFIXES = {
    ".txt", ".md", ".rst", ".py", ".js", ".ts", ".tsx", ".jsx", ".json", ".yaml", ".yml",
    ".toml", ".ini", ".cfg", ".csv", ".sql", ".html", ".css", ".xml", ".log", ".ps1", ".sh",
}
CODE_SUFFIXES = {".py", ".js", ".ts", ".tsx", ".jsx", ".sql", ".ps1", ".sh", ".html", ".css"}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
AUDIO_SUFFIXES = {".wav", ".mp3", ".m4a", ".flac", ".ogg", ".aac"}
VIDEO_SUFFIXES = {".mp4", ".mov", ".mkv", ".webm", ".avi"}


@dataclass(frozen=True)
class ExtractedAsset:
    asset: MultimodalAsset
    evidence: tuple[MultimodalEvidence, ...]


class WorkspaceAssetResolver:
    """Resolve files inside a governed workspace and reject path escape."""

    def __init__(self, workspace: Path | str):
        self.workspace = Path(workspace).resolve()

    def resolve(self, raw_path: str) -> Path:
        if not str(raw_path or "").strip():
            raise MultimodalRuntimeError("asset path is required")
        candidate = Path(raw_path)
        if not candidate.is_absolute():
            candidate = self.workspace / candidate
        candidate = candidate.resolve()
        try:
            candidate.relative_to(self.workspace)
        except ValueError as exc:
            raise MultimodalRuntimeError("asset path must remain inside the workspace") from exc
        if not candidate.is_file():
            raise MultimodalRuntimeError(f"asset file does not exist: {raw_path}")
        return candidate

    def relative(self, path: Path) -> str:
        return path.resolve().relative_to(self.workspace).as_posix()


class MultimodalExtractor:
    """Native-first extraction adapters for v10 multimodal assets.

    Native document/table/code/transcript extraction is preferred. Images, raw audio,
    and video are inventoried with cryptographic provenance but are not falsely
    interpreted without a capable provider. PDF pages are emitted individually so
    page provenance is never lost during downstream reasoning.
    """

    def __init__(self, workspace: Path | str):
        self.resolver = WorkspaceAssetResolver(workspace)

    @staticmethod
    def _sha256(path: Path) -> str:
        digest = hashlib.sha256()
        with path.open("rb") as handle:
            for block in iter(lambda: handle.read(1024 * 1024), b""):
                digest.update(block)
        return digest.hexdigest()

    @staticmethod
    def _evidence(asset_id: str, kind: EvidenceKind, content: str, suffix: str) -> MultimodalEvidence:
        return MultimodalEvidence(
            evidence_id=f"{asset_id}:{suffix}",
            asset_id=asset_id,
            kind=kind,
            content=content,
            provider="native_extractor",
            model="deterministic-parser",
            confidence=1.0,
        )

    def extract(self, raw_path: str, *, asset_id: str | None = None) -> tuple[ExtractedAsset, ...]:
        path = self.resolver.resolve(raw_path)
        suffix = path.suffix.lower()
        base_id = str(asset_id or path.stem).strip()
        if not base_id:
            raise MultimodalRuntimeError("asset id is required")
        if suffix == ".pdf":
            return self._extract_pdf(path, base_id)
        if suffix == ".docx":
            return (self._extract_docx(path, base_id),)
        if suffix == ".xlsx":
            return (self._extract_xlsx(path, base_id),)
        if suffix == ".pptx":
            return (self._extract_pptx(path, base_id),)
        if suffix in TEXT_SUFFIXES or suffix == "":
            return (self._extract_text(path, base_id),)
        if suffix in IMAGE_SUFFIXES:
            return (self._inventory_binary(path, base_id, Modality.IMAGE),)
        if suffix in AUDIO_SUFFIXES:
            return (self._inventory_binary(path, base_id, Modality.AUDIO),)
        if suffix in VIDEO_SUFFIXES:
            return (self._inventory_binary(path, base_id, Modality.VIDEO),)
        raise MultimodalRuntimeError(f"unsupported multimodal file type: {suffix or '<none>'}")

    def _base_asset(self, path: Path, asset_id: str, modality: Modality, *, native_text: str = "", page: int | None = None) -> MultimodalAsset:
        return MultimodalAsset(
            asset_id=asset_id,
            modality=modality,
            source_ref=self.resolver.relative(path),
            sha256=self._sha256(path),
            mime_type="",
            page=page,
            native_text=native_text,
        )

    def _extract_text(self, path: Path, asset_id: str) -> ExtractedAsset:
        text = path.read_text(encoding="utf-8", errors="replace")
        modality = Modality.CODE if path.suffix.lower() in CODE_SUFFIXES else Modality.TEXT
        if path.suffix.lower() == ".csv":
            rows = []
            with path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
                for index, row in enumerate(csv.reader(handle)):
                    if index >= 10000:
                        break
                    rows.append(" | ".join(row))
            text = "\n".join(rows)
            modality = Modality.SPREADSHEET
        asset = self._base_asset(path, asset_id, modality, native_text=text)
        kind = EvidenceKind.CODE if modality == Modality.CODE else (EvidenceKind.TABLE if modality == Modality.SPREADSHEET else EvidenceKind.NATIVE_TEXT)
        evidence = () if not text.strip() else (self._evidence(asset_id, kind, text, "native"),)
        return ExtractedAsset(asset, evidence)

    def _extract_pdf(self, path: Path, base_id: str) -> tuple[ExtractedAsset, ...]:
        reader = PdfReader(str(path))
        outputs: list[ExtractedAsset] = []
        file_hash = self._sha256(path)
        for page_number, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            asset_id = f"{base_id}:page:{page_number}"
            asset = MultimodalAsset(
                asset_id=asset_id,
                modality=Modality.PDF,
                source_ref=self.resolver.relative(path),
                sha256=file_hash,
                page=page_number,
                native_text=text,
            )
            evidence = () if not text.strip() else (self._evidence(asset_id, EvidenceKind.NATIVE_TEXT, text, "native"),)
            outputs.append(ExtractedAsset(asset, evidence))
        if not outputs:
            raise MultimodalRuntimeError("PDF contains no pages")
        return tuple(outputs)

    def _extract_docx(self, path: Path, asset_id: str) -> ExtractedAsset:
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text]
        tables: list[str] = []
        for table in doc.tables:
            for row in table.rows:
                tables.append(" | ".join(cell.text for cell in row.cells))
        text = "\n".join(paragraphs + tables)
        asset = self._base_asset(path, asset_id, Modality.DOCUMENT, native_text=text)
        evidence: list[MultimodalEvidence] = []
        if paragraphs:
            evidence.append(self._evidence(asset_id, EvidenceKind.NATIVE_TEXT, "\n".join(paragraphs), "text"))
        if tables:
            evidence.append(self._evidence(asset_id, EvidenceKind.TABLE, "\n".join(tables), "tables"))
        return ExtractedAsset(asset, tuple(evidence))

    def _extract_xlsx(self, path: Path, asset_id: str) -> ExtractedAsset:
        wb = load_workbook(str(path), read_only=True, data_only=True)
        try:
            sheets: list[dict[str, object]] = []
            rendered: list[str] = []
            for ws in wb.worksheets:
                rows: list[list[str]] = []
                for index, row in enumerate(ws.iter_rows(values_only=True)):
                    if index >= 10000:
                        break
                    values = ["" if value is None else str(value) for value in row]
                    rows.append(values)
                    rendered.append(" | ".join(values))
                sheets.append({"title": ws.title, "rows": rows})
            text = "\n".join(rendered)
            asset = self._base_asset(path, asset_id, Modality.SPREADSHEET, native_text=text)
            structure = json.dumps({"sheets": sheets}, ensure_ascii=False)
            evidence = (
                self._evidence(asset_id, EvidenceKind.STRUCTURE, structure, "structure"),
                self._evidence(asset_id, EvidenceKind.TABLE, text or "workbook contains no populated cells", "table"),
            )
            return ExtractedAsset(asset, evidence)
        finally:
            wb.close()

    def _extract_pptx(self, path: Path, asset_id: str) -> ExtractedAsset:
        prs = Presentation(str(path))
        slides: list[dict[str, object]] = []
        rendered: list[str] = []
        for index, slide in enumerate(prs.slides, start=1):
            texts = [str(shape.text) for shape in slide.shapes if hasattr(shape, "text") and str(shape.text).strip()]
            slides.append({"slide": index, "text": texts})
            rendered.append(f"# Slide {index}\n" + "\n".join(texts))
        text = "\n\n".join(rendered)
        asset = self._base_asset(path, asset_id, Modality.PRESENTATION, native_text=text)
        evidence: list[MultimodalEvidence] = [self._evidence(asset_id, EvidenceKind.STRUCTURE, json.dumps({"slides": slides}, ensure_ascii=False), "structure")]
        if text.strip():
            evidence.append(self._evidence(asset_id, EvidenceKind.NATIVE_TEXT, text, "text"))
        return ExtractedAsset(asset, tuple(evidence))

    def _inventory_binary(self, path: Path, asset_id: str, modality: Modality) -> ExtractedAsset:
        stat = path.stat()
        asset = self._base_asset(path, asset_id, modality)
        metadata = json.dumps({"filename": path.name, "size_bytes": stat.st_size, "suffix": path.suffix.lower()}, sort_keys=True)
        evidence = (self._evidence(asset_id, EvidenceKind.METADATA, metadata, "metadata"),)
        return ExtractedAsset(asset, evidence)

    @staticmethod
    def attach_transcript(asset: MultimodalAsset, transcript: str, *, provider: str, model: str) -> MultimodalEvidence:
        asset.validate()
        if asset.modality not in {Modality.AUDIO, Modality.VIDEO}:
            raise MultimodalRuntimeError("transcripts may only be attached to audio or video assets")
        if not transcript.strip() or not provider.strip() or not model.strip():
            raise MultimodalRuntimeError("transcript evidence requires text, provider, and model")
        return MultimodalEvidence(
            evidence_id=f"{asset.asset_id}:transcript",
            asset_id=asset.asset_id,
            kind=EvidenceKind.TRANSCRIPT,
            content=transcript.strip(),
            provider=provider.strip(),
            model=model.strip(),
        )


__all__ = ["ExtractedAsset", "MultimodalExtractor", "WorkspaceAssetResolver"]
