from __future__ import annotations

import csv
import json
import re
from pathlib import Path
from typing import Any, Iterable

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from app.db import Database, utc_now
from app.tools.filesystem import TEXT_EXTENSIONS, WorkspaceFS


SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | {".pdf", ".docx", ".xlsx", ".pptx"}


def chunk_text(text: str, size: int = 1800, overlap: int = 250) -> list[str]:
    cleaned = re.sub(r"\r\n?", "\n", text)
    cleaned = re.sub(r"[ \t]+", " ", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned).strip()
    if not cleaned:
        return []
    chunks: list[str] = []
    start = 0
    while start < len(cleaned):
        end = min(start + size, len(cleaned))
        if end < len(cleaned):
            boundary = max(cleaned.rfind("\n\n", start, end), cleaned.rfind(". ", start, end))
            if boundary > start + size // 2:
                end = boundary + 1
        chunks.append(cleaned[start:end].strip())
        if end >= len(cleaned):
            break
        start = max(end - overlap, start + 1)
    return chunks


class KnowledgeBase:
    def __init__(self, db: Database, fs: WorkspaceFS):
        self.db = db
        self.fs = fs

    def extract_text(self, path: Path) -> str:
        suffix = path.suffix.lower()
        if suffix in TEXT_EXTENSIONS or suffix == "":
            return path.read_text(encoding="utf-8", errors="replace")
        if suffix == ".pdf":
            reader = PdfReader(str(path))
            return "\n\n".join(page.extract_text() or "" for page in reader.pages)
        if suffix == ".docx":
            doc = Document(str(path))
            parts = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
            return "\n".join(parts)
        if suffix == ".xlsx":
            wb = load_workbook(str(path), read_only=True, data_only=True)
            parts: list[str] = []
            for ws in wb.worksheets:
                parts.append(f"# Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    parts.append(" | ".join("" if value is None else str(value) for value in row))
            wb.close()
            return "\n".join(parts)
        if suffix == ".pptx":
            prs = Presentation(str(path))
            parts: list[str] = []
            for index, slide in enumerate(prs.slides, start=1):
                parts.append(f"# Slide {index}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            return "\n".join(parts)
        raise ValueError(f"Unsupported file type: {suffix}")

    def _upsert_document(self, path: Path, text: str) -> int:
        rel = self.fs.relative(path)
        stat = path.stat()
        chunks = chunk_text(text)
        with self.db.connect() as db:
            row = db.execute("SELECT id FROM knowledge_documents WHERE path = ?", (rel,)).fetchone()
            if row:
                document_id = int(row["id"])
                old_chunk_ids = [r["id"] for r in db.execute("SELECT id FROM knowledge_chunks WHERE document_id = ?", (document_id,)).fetchall()]
                if old_chunk_ids:
                    placeholders = ",".join("?" for _ in old_chunk_ids)
                    db.execute(f"DELETE FROM knowledge_fts WHERE chunk_id IN ({placeholders})", tuple(old_chunk_ids))
                db.execute("DELETE FROM knowledge_chunks WHERE document_id = ?", (document_id,))
                db.execute(
                    "UPDATE knowledge_documents SET modified_ns=?, size_bytes=?, indexed_at=? WHERE id=?",
                    (stat.st_mtime_ns, stat.st_size, utc_now(), document_id),
                )
            else:
                cursor = db.execute(
                    "INSERT INTO knowledge_documents(path, modified_ns, size_bytes, indexed_at) VALUES (?, ?, ?, ?)",
                    (rel, stat.st_mtime_ns, stat.st_size, utc_now()),
                )
                document_id = int(cursor.lastrowid)
            for index, content in enumerate(chunks):
                cursor = db.execute(
                    "INSERT INTO knowledge_chunks(document_id, chunk_index, content) VALUES (?, ?, ?)",
                    (document_id, index, content),
                )
                chunk_id = int(cursor.lastrowid)
                db.execute(
                    "INSERT INTO knowledge_fts(content, path, chunk_id) VALUES (?, ?, ?)",
                    (content, rel, chunk_id),
                )
        return len(chunks)

    def index_workspace(self, path: str = ".", force: bool = False, max_files: int = 1000) -> dict[str, Any]:
        base = self.fs.resolve(path)
        candidates: Iterable[Path] = [base] if base.is_file() else base.rglob("*")
        indexed = 0
        skipped = 0
        failed: list[dict[str, str]] = []
        chunk_count = 0
        for candidate in candidates:
            if indexed + skipped >= max_files:
                break
            if not candidate.is_file() or candidate.suffix.lower() not in SUPPORTED_EXTENSIONS:
                continue
            if any(part in {".git", ".venv", "node_modules", "__pycache__"} for part in candidate.parts):
                continue
            rel = self.fs.relative(candidate)
            stat = candidate.stat()
            with self.db.connect() as db:
                row = db.execute(
                    "SELECT modified_ns, size_bytes FROM knowledge_documents WHERE path = ?",
                    (rel,),
                ).fetchone()
            if row and not force and row["modified_ns"] == stat.st_mtime_ns and row["size_bytes"] == stat.st_size:
                skipped += 1
                continue
            try:
                text = self.extract_text(candidate)
                chunk_count += self._upsert_document(candidate, text)
                indexed += 1
            except Exception as exc:  # noqa: BLE001
                failed.append({"path": rel, "error": str(exc)})
        return {"ok": True, "indexed": indexed, "skipped": skipped, "chunks": chunk_count, "failed": failed[:25]}

    def search(self, query: str, limit: int = 8) -> dict[str, Any]:
        tokens = re.findall(r"[A-Za-z0-9_\-]{2,}", query)
        if not tokens:
            return {"ok": True, "results": []}
        fts_query = " OR ".join(f'"{token}"' for token in tokens[:16])
        with self.db.connect() as db:
            try:
                rows = db.execute(
                    """
                    SELECT path, content, bm25(knowledge_fts) AS score
                    FROM knowledge_fts
                    WHERE knowledge_fts MATCH ?
                    ORDER BY score
                    LIMIT ?
                    """,
                    (fts_query, max(1, min(limit, 30))),
                ).fetchall()
            except Exception:
                rows = []
        return {
            "ok": True,
            "results": [
                {"path": row["path"], "content": row["content"], "score": row["score"]}
                for row in rows
            ],
        }

    def stats(self) -> dict[str, Any]:
        with self.db.connect() as db:
            documents = db.execute("SELECT COUNT(*) AS count FROM knowledge_documents").fetchone()["count"]
            chunks = db.execute("SELECT COUNT(*) AS count FROM knowledge_chunks").fetchone()["count"]
        return {"documents": documents, "chunks": chunks}